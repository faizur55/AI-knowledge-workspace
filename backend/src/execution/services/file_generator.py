"""
File Generator

Generates output files in various formats (PDF, DOCX, PPTX, XLSX, CSV, etc.)
"""

import os
import io
import json
import csv
from typing import Dict, Any, Optional, List, Union
from datetime import datetime
import base64

from src.core.logging import logger


class FileGenerator:
    """
    Generator for various file formats.
    
    Supports:
    - PDF
    - DOCX
    - PPTX
    - XLSX
    - CSV
    - JSON
    - Markdown
    - Text
    - HTML
    """

    def __init__(self, output_dir: str = "outputs"):
        """Initialize the file generator."""
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    async def generate(
        self,
        format: str,
        data: Dict[str, Any],
        filename: str,
        execution_id: str
    ) -> Dict[str, Any]:
        """
        Generate a file in the specified format.
        
        Args:
            format: Output format (pdf, docx, pptx, xlsx, csv, json, markdown, text, html)
            data: Data to include in the file
            filename: Output filename (without extension)
            execution_id: Execution ID for organizing files
            
        Returns:
            File information dict
        """
        # Create execution directory
        exec_dir = os.path.join(self.output_dir, execution_id)
        os.makedirs(exec_dir, exist_ok=True)
        
        # Generate file based on format
        generators = {
            "pdf": self._generate_pdf,
            "docx": self._generate_docx,
            "pptx": self._generate_pptx,
            "xlsx": self._generate_xlsx,
            "csv": self._generate_csv,
            "json": self._generate_json,
            "markdown": self._generate_markdown,
            "text": self._generate_text,
            "html": self._generate_html,
        }
        
        generator = generators.get(format.lower())
        if not generator:
            raise ValueError(f"Unsupported format: {format}")
        
        result = await generator(data, filename, exec_dir)
        return result

    async def _generate_pdf(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate PDF file."""
        # Use reportlab or weasyprint for PDF generation
        # For now, generate HTML and convert (or use simple approach)
        
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.enums import TA_LEFT, TA_CENTER
            
            filepath = os.path.join(output_dir, f"{filename}.pdf")
            
            doc = SimpleDocTemplate(
                filepath,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18
            )
            
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=30,
                alignment=TA_CENTER
            )
            
            story = []
            
            # Title
            if data.get("title"):
                story.append(Paragraph(data["title"], title_style))
                story.append(Spacer(1, 0.2 * inch))
            
            # Content
            content = data.get("content", "")
            if isinstance(content, list):
                for item in content:
                    story.append(Paragraph(str(item), styles['Normal']))
                    story.append(Spacer(1, 0.1 * inch))
            else:
                story.append(Paragraph(content, styles['Normal']))
            
            doc.build(story)
            
            file_size = os.path.getsize(filepath)
            
            return {
                "filename": f"{filename}.pdf",
                "format": "pdf",
                "path": filepath,
                "size_bytes": file_size,
                "content_type": "application/pdf",
                "success": True
            }
            
        except ImportError:
            # Fallback: generate HTML that can be printed to PDF
            html_data = {
                "title": data.get("title", "Report"),
                "content": data.get("content", ""),
                "generated_at": datetime.utcnow().isoformat()
            }
            return await self._generate_html(html_data, filename, output_dir)

    async def _generate_docx(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate DOCX file."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
            
            filepath = os.path.join(output_dir, f"{filename}.docx")
            
            doc = Document()
            
            # Title
            if data.get("title"):
                title = doc.add_heading(data["title"], 0)
                title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            
            # Sections
            sections = data.get("sections", [])
            if not sections and data.get("content"):
                sections = [{"title": "", "content": data["content"]}]
            
            for section in sections:
                if section.get("title"):
                    doc.add_heading(section["title"], level=1)
                
                content = section.get("content", "")
                if isinstance(content, list):
                    for item in content:
                        doc.add_paragraph(str(item))
                elif content:
                    doc.add_paragraph(str(content))
                
                doc.add_paragraph()  # Spacer
            
            # Tables
            tables = data.get("tables", [])
            for table_data in tables:
                if table_data.get("headers") and table_data.get("rows"):
                    table = doc.add_table(rows=1, cols=len(table_data["headers"]))
                    table.style = 'Light Grid Accent 1'
                    
                    # Header row
                    hdr_cells = table.rows[0].cells
                    for i, header in enumerate(table_data["headers"]):
                        hdr_cells[i].text = str(header)
                    
                    # Data rows
                    for row in table_data["rows"]:
                        row_cells = table.add_row().cells
                        for i, cell in enumerate(row):
                            row_cells[i].text = str(cell)
                
                doc.add_paragraph()
            
            doc.save(filepath)
            file_size = os.path.getsize(filepath)
            
            return {
                "filename": f"{filename}.docx",
                "format": "docx",
                "path": filepath,
                "size_bytes": file_size,
                "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "success": True
            }
            
        except ImportError:
            # Fallback: generate markdown
            md_data = {
                "title": data.get("title", ""),
                "content": data.get("content", ""),
                "sections": data.get("sections", [])
            }
            return await self._generate_markdown(md_data, filename, output_dir)

    async def _generate_pptx(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate PPTX file."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RgbColor
            
            filepath = os.path.join(output_dir, f"{filename}.pptx")
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            if data.get("title"):
                title = slide.shapes.title
                title.text = data["title"]
            
            if data.get("subtitle"):
                subtitle = slide.placeholders[1]
                subtitle.text = data["subtitle"]
            
            # Content slides
            slides_data = data.get("slides", [])
            if not slides_data and data.get("content"):
                slides_data = [{"title": "Overview", "content": data["content"]}]
            
            for slide_data in slides_data:
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                if slide_data.get("title"):
                    slide.shapes.title.text = slide_data["title"]
                
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                
                content = slide_data.get("content", "")
                if isinstance(content, list):
                    for i, item in enumerate(content):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = str(item)
                        p.level = slide_data.get("level", 0)
                elif content:
                    tf.paragraphs[0].text = str(content)
            
            prs.save(filepath)
            file_size = os.path.getsize(filepath)
            
            return {
                "filename": f"{filename}.pptx",
                "format": "pptx",
                "path": filepath,
                "size_bytes": file_size,
                "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "success": True
            }
            
        except ImportError:
            # Fallback: return markdown
            return await self._generate_markdown(data, filename, output_dir)

    async def _generate_xlsx(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate XLSX file."""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill
            from openpyxl.utils import get_column_letter
            
            filepath = os.path.join(output_dir, f"{filename}.xlsx")
            
            wb = Workbook()
            
            # Remove default sheet
            if 'Sheet' in wb.sheetnames:
                del wb['Sheet']
            
            # Create sheets from data
            sheets = data.get("sheets", [])
            
            if not sheets and data.get("headers") and data.get("rows"):
                sheets = [{"name": "Data", "headers": data["headers"], "rows": data["rows"]}]
            
            for sheet_data in sheets:
                sheet_name = sheet_data.get("name", "Sheet")
                ws = wb.create_sheet(title=sheet_name[:31])  # Excel limit
                
                headers = sheet_data.get("headers", [])
                rows = sheet_data.get("rows", [])
                
                # Style for headers
                header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
                header_font = Font(bold=True, color="FFFFFF")
                
                # Add headers
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col)
                    cell.value = header
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = Alignment(horizontal="center")
                
                # Add data rows
                for row_idx, row_data in enumerate(rows, 2):
                    if isinstance(row_data, dict):
                        for col, header in enumerate(headers, 1):
                            ws.cell(row=row_idx, column=col).value = row_data.get(header, "")
                    elif isinstance(row_data, list):
                        for col, value in enumerate(row_data, 1):
                            ws.cell(row=row_idx, column=col).value = value
                
                # Auto-adjust column widths
                for col in range(1, len(headers) + 1):
                    ws.column_dimensions[get_column_letter(col)].width = 15
            
            wb.save(filepath)
            file_size = os.path.getsize(filepath)
            
            return {
                "filename": f"{filename}.xlsx",
                "format": "xlsx",
                "path": filepath,
                "size_bytes": file_size,
                "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "success": True
            }
            
        except ImportError:
            # Fallback: generate CSV
            if data.get("headers") and data.get("rows"):
                csv_data = {"headers": data["headers"], "rows": data["rows"]}
                return await self._generate_csv(csv_data, filename, output_dir)
            return await self._generate_json(data, filename, output_dir)

    async def _generate_csv(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate CSV file."""
        filepath = os.path.join(output_dir, f"{filename}.csv")
        
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        
        with open(filepath, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            
            if headers:
                writer.writerow(headers)
            
            for row in rows:
                if isinstance(row, dict):
                    writer.writerow([row.get(h, "") for h in headers])
                else:
                    writer.writerow(row)
        
        file_size = os.path.getsize(filepath)
        
        return {
            "filename": f"{filename}.csv",
            "format": "csv",
            "path": filepath,
            "size_bytes": file_size,
            "content_type": "text/csv",
            "success": True
        }

    async def _generate_json(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate JSON file."""
        filepath = os.path.join(output_dir, f"{filename}.json")
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        
        file_size = os.path.getsize(filepath)
        
        return {
            "filename": f"{filename}.json",
            "format": "json",
            "path": filepath,
            "size_bytes": file_size,
            "content_type": "application/json",
            "success": True
        }

    async def _generate_markdown(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate Markdown file."""
        filepath = os.path.join(output_dir, f"{filename}.md")
        
        lines = []
        
        # Title
        if data.get("title"):
            lines.append(f"# {data['title']}")
            lines.append("")
        
        # Sections
        sections = data.get("sections", [])
        if not sections and data.get("content"):
            sections = [{"title": "", "content": data["content"]}]
        
        for section in sections:
            if section.get("title"):
                lines.append(f"## {section['title']}")
                lines.append("")
            
            content = section.get("content", "")
            if isinstance(content, list):
                for item in content:
                    lines.append(f"- {item}")
            elif content:
                lines.append(content)
            lines.append("")
        
        # Tables
        tables = data.get("tables", [])
        for table in tables:
            if table.get("headers") and table.get("rows"):
                lines.append(f"### {table.get('title', 'Table')}")
                lines.append("")
                lines.append("| " + " | ".join(str(h) for h in table["headers"]) + " |")
                lines.append("| " + " | ".join("---" for _ in table["headers"]) + " |")
                for row in table["rows"]:
                    lines.append("| " + " | ".join(str(c) for c in row) + " |")
                lines.append("")
        
        # Metadata
        lines.append("---")
        lines.append(f"*Generated: {datetime.utcnow().isoformat()}*")
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("\n".join(lines))
        
        file_size = os.path.getsize(filepath)
        
        return {
            "filename": f"{filename}.md",
            "format": "markdown",
            "path": filepath,
            "size_bytes": file_size,
            "content_type": "text/markdown",
            "success": True
        }

    async def _generate_text(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate plain text file."""
        filepath = os.path.join(output_dir, f"{filename}.txt")
        
        lines = []
        
        if data.get("title"):
            lines.append(data["title"].upper())
            lines.append("=" * len(data["title"]))
            lines.append("")
        
        content = data.get("content", "")
        if isinstance(content, list):
            for item in content:
                lines.append(f"• {item}")
        elif content:
            lines.append(content)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("\n".join(lines))
        
        file_size = os.path.getsize(filepath)
        
        return {
            "filename": f"{filename}.txt",
            "format": "text",
            "path": filepath,
            "size_bytes": file_size,
            "content_type": "text/plain",
            "success": True
        }

    async def _generate_html(
        self,
        data: Dict[str, Any],
        filename: str,
        output_dir: str
    ) -> Dict[str, Any]:
        """Generate HTML file."""
        filepath = os.path.join(output_dir, f"{filename}.html")
        
        html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{data.get('title', 'Document')}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
            line-height: 1.6;
            max-width: 800px;
            margin: 0 auto;
            padding: 40px 20px;
            color: #333;
        }}
        h1 {{ color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 10px; }}
        h2 {{ color: #34495e; margin-top: 30px; }}
        table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
        th, td {{ border: 1px solid #ddd; padding: 12px; text-align: left; }}
        th {{ background-color: #3498db; color: white; }}
        tr:nth-child(even) {{ background-color: #f9f9f9; }}
        .meta {{ color: #7f8c8d; font-size: 0.9em; margin-top: 40px; }}
    </style>
</head>
<body>
    <h1>{data.get('title', 'Document')}</h1>
"""
        
        content = data.get("content", "")
        if isinstance(content, list):
            html_content += "    <ul>\n"
            for item in content:
                html_content += f"        <li>{item}</li>\n"
            html_content += "    </ul>\n"
        elif content:
            html_content += f"    <p>{content}</p>\n"
        
        html_content += f"""
    <div class="meta">
        <p>Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}</p>
    </div>
</body>
</html>"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        file_size = os.path.getsize(filepath)
        
        return {
            "filename": f"{filename}.html",
            "format": "html",
            "path": filepath,
            "size_bytes": file_size,
            "content_type": "text/html",
            "success": True
        }

    def generate_base64(self, filepath: str) -> str:
        """Generate base64 encoded version of a file."""
        with open(filepath, 'rb') as f:
            return base64.b64encode(f.read()).decode('utf-8')

    def get_download_url(self, filepath: str) -> str:
        """Get download URL for a file."""
        # In production, this would generate a signed URL
        return f"/outputs/{os.path.basename(os.path.dirname(filepath))}/{os.path.basename(filepath)}"
